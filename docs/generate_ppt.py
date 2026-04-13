#!/usr/bin/env python3
"""DataAgent 数据智能体平台 — 战略规划 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 配色
# ============================================================
C_DARK      = RGBColor(0x0F, 0x17, 0x2A)
C_BLUE      = RGBColor(0x3B, 0x82, 0xF6)
C_PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
C_GREEN     = RGBColor(0x05, 0x96, 0x69)
C_AMBER     = RGBColor(0xD9, 0x77, 0x06)
C_RED       = RGBColor(0xDC, 0x26, 0x26)
C_INDIGO    = RGBColor(0x4F, 0x46, 0xE5)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xF5, 0xF7, 0xFA)
C_GRAY      = RGBColor(0x86, 0x90, 0x9C)
C_DARKGRAY  = RGBColor(0x4E, 0x59, 0x69)
C_TEXT      = RGBColor(0x1D, 0x21, 0x29)

# 阶段颜色
PHASE_COLORS = [C_BLUE, C_PURPLE, C_GREEN, C_AMBER]
PHASE_NAMES  = ["智能SQL开发中心", "智能分析中心", "数据治理中心", "数据管道编排中心"]
PHASE_QUARTERS = ["2026 Q2-Q3", "2026 Q4-2027 Q1", "2027 Q1-Q2", "2027 Q2-Q3"]
CATALOG_NAMES = ["技术元数据目录", "业务语义目录", "智能治理目录", "全景数据图谱"]

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ============================================================
# 工具函数
# ============================================================
def add_rect(slide, left, top, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def set_text(shape, text, font_size=14, color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_text_box(slide, left, top, w, h, text, font_size=14, color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    return tf

def add_bullet_list(tf, items, font_size=12, color=C_DARKGRAY, bullet_color=C_GRAY):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)
        p.level = 0

def slide_bg(slide, color=C_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_section_header(slide, title, subtitle=""):
    # 左侧蓝色竖条
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SH, C_BLUE)
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(10), Inches(0.6), title, 28, C_DARK, True)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(0.95), Inches(10), Inches(0.4), subtitle, 14, C_GRAY)

def add_card(slide, left, top, w, h, title, items, header_color, title_size=13):
    card = add_rounded_rect(slide, left, top, w, h, C_WHITE, RGBColor(0xE2, 0xE8, 0xF0))
    # header bar
    hdr = add_rect(slide, left + Inches(0.05), top + Inches(0.05), w - Inches(0.1), Inches(0.4), header_color)
    hdr.shape_type  # rounded rect clips; just overlay a rect
    set_text(hdr, title, title_size, C_WHITE, True, PP_ALIGN.CENTER)
    # body
    body_top = top + Inches(0.5)
    tf = add_text_box(slide, left + Inches(0.15), body_top, w - Inches(0.3), h - Inches(0.6), "", 11, C_DARKGRAY)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "· " + item
        p.font.size = Pt(11)
        p.font.color.rgb = C_DARKGRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(3)

# ============================================================
# P1 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(s, C_DARK)
add_text_box(s, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6), "DATAAGENT PLATFORM", 20, C_GRAY, False, PP_ALIGN.CENTER)
add_text_box(s, Inches(1.5), Inches(2.3), Inches(10), Inches(1), "DataAgent 数据智能体平台", 44, C_WHITE, True, PP_ALIGN.CENTER)
# 蓝色分割线
add_rect(s, Inches(5.6), Inches(3.5), Inches(2), Inches(0.04), C_BLUE)
add_text_box(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6), "战略规划蓝图", 22, RGBColor(0x94,0xA3,0xB8), False, PP_ALIGN.CENTER)
add_text_box(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4), "2026年4月  |  V1.0", 14, RGBColor(0x64,0x74,0x8B), False, PP_ALIGN.CENTER)

# ============================================================
# P2 项目背景
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "项目背景", "数据领域面临三组结构性矛盾")

problems = [
    ("需求增长 vs 供给不足", "数据需求爆发式增长，但开发高度依赖少数工程师，需求排期积压，平均交付周期3-5天", C_BLUE),
    ("资产膨胀 vs 治理滞后", "数据资产规模指数级增长，但治理仍靠人工，表结构文档过时、血缘缺失、敏感数据分布不明", C_PURPLE),
    ("技术进步 vs 体验脱节", "已建成先进Lakehouse架构，但使用门槛高——业务人员必须掌握SQL才能获取数据", C_GREEN),
]

for i, (title, desc, color) in enumerate(problems):
    left = Inches(0.5) + i * Inches(4.1)
    top = Inches(1.8)
    card = add_rounded_rect(s, left, top, Inches(3.9), Inches(2.6), C_WHITE)
    add_rect(s, left, top, Inches(3.9), Inches(0.08), color)
    add_text_box(s, left + Inches(0.3), top + Inches(0.3), Inches(3.3), Inches(0.5), title, 16, color, True)
    add_text_box(s, left + Inches(0.3), top + Inches(0.9), Inches(3.3), Inches(1.5), desc, 12, C_DARKGRAY)

# 底部总结
add_text_box(s, Inches(0.5), Inches(5.0), Inches(12), Inches(1.2),
    '大语言模型技术的成熟为系统性解决上述矛盾提供了可能。DataAgent以AI为内核，重新定义数据工作范式，将数据能力从"专家垄断"转变为"全员赋能"。',
    14, C_DARKGRAY, False, PP_ALIGN.LEFT)

# ============================================================
# P3 建设目标
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "建设目标", "四大核心目标，可量化、可考核")

goals = [
    ("提升数据开发效率", "NL2SQL + 智能优化", "需求交付周期\n天级 → 分钟级", C_BLUE),
    ("实现数据自助分析", "对话式分析 + 智能可视化", "业务自助分析占比\n≥ 60%", C_PURPLE),
    ("构建智能治理体系", "血缘追踪 + 质量守护", "核心资产入目入册\n100%", C_GREEN),
    ("自动化数据管道", "自然语言管道设计", "管道开发效率\n提升 5 倍", C_AMBER),
]

for i, (title, sub, kpi, color) in enumerate(goals):
    left = Inches(0.5) + i * Inches(3.1)
    top = Inches(1.8)
    card = add_rounded_rect(s, left, top, Inches(2.9), Inches(4.2), C_WHITE)
    add_rect(s, left, top, Inches(2.9), Inches(0.08), color)
    # 数字编号
    num = add_rounded_rect(s, left + Inches(1.05), top + Inches(0.3), Inches(0.7), Inches(0.5), color)
    set_text(num, f"0{i+1}", 16, C_WHITE, True, PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.2), top + Inches(1.0), Inches(2.5), Inches(0.5), title, 16, C_DARK, True, PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.2), top + Inches(1.5), Inches(2.5), Inches(0.5), sub, 11, C_GRAY, False, PP_ALIGN.CENTER)
    # KPI
    kpi_box = add_rounded_rect(s, left + Inches(0.3), top + Inches(2.4), Inches(2.3), Inches(1.3), RGBColor(0xF8,0xFA,0xFC))
    tf = set_text(kpi_box, kpi, 15, color, True, PP_ALIGN.CENTER)
    tf.paragraphs[0].space_after = Pt(4)
    tf.word_wrap = True
    # vertical center
    kpi_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# P4 总体架构
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_WHITE)
add_section_header(s, "平台总体架构", "三层分离 · Agent驱动 · 目录贯穿")

layers = [
    ("应用服务层", ["智能SQL开发中心", "智能分析中心", "数据治理中心", "管道编排中心"], C_BLUE),
    ("服务网关层", ["认证鉴权", "请求路由", "会话管理", "限流熔断"], C_PURPLE),
    ("AI 引擎层", ["Router Agent", "SQL Agent", "Analytics Agent", "Governance Agent", "Pipeline Agent"], C_GREEN),
    ("基础设施连接层", ["Trino", "Spark", "Iceberg", "Pulsar", "Airflow"], C_AMBER),
    ("数据与知识存储层", ["PostgreSQL", "Milvus", "Redis", "S3"], C_RED),
]

y_start = Inches(1.5)
layer_h = Inches(0.85)
gap = Inches(0.08)

for i, (name, items, color) in enumerate(layers):
    top = y_start + i * (layer_h + gap)
    # 左侧标签
    label = add_rounded_rect(s, Inches(0.3), top, Inches(2.0), layer_h, color)
    set_text(label, name, 12, C_WHITE, True, PP_ALIGN.CENTER)
    label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 右侧内容
    for j, item in enumerate(items):
        bx = Inches(2.5) + j * Inches(2.1)
        box = add_rounded_rect(s, bx, top + Inches(0.1), Inches(1.95), layer_h - Inches(0.2), RGBColor(0xF8,0xFA,0xFC), RGBColor(0xE2,0xE8,0xF0))
        set_text(box, item, 11, C_DARKGRAY, False, PP_ALIGN.CENTER)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 底部：数据目录贯穿 + 两条保障线
catalog_bar = add_rounded_rect(s, Inches(0.3), Inches(5.9), Inches(9.5), Inches(0.55), RGBColor(0xEF,0xF6,0xFF), C_BLUE)
tf = set_text(catalog_bar, "📖 数据目录（Data Catalog）— 贯穿全平台的数据资产基座：技术元数据 → 业务语义 → 智能治理 → 全景图谱", 11, RGBColor(0x1E,0x40,0xAF), True, PP_ALIGN.CENTER)

sec_bar = add_rounded_rect(s, Inches(0.3), Inches(6.55), Inches(4.6), Inches(0.45), RGBColor(0xFE,0xF2,0xF2))
set_text(sec_bar, "🔒 安全合规保障 — 认证授权 · 数据脱敏 · SQL审查 · 审计日志", 10, RGBColor(0x99,0x1B,0x1B), False, PP_ALIGN.CENTER)

know_bar = add_rounded_rect(s, Inches(5.1), Inches(6.55), Inches(4.7), Inches(0.45), RGBColor(0xF0,0xFD,0xF4))
set_text(know_bar, "📚 知识管理保障 — 元数据采集 · RAG检索 · 反馈闭环 · 知识生命周期", 10, RGBColor(0x16,0x65,0x34), False, PP_ALIGN.CENTER)

# 右侧技术标注
tech_box = add_rounded_rect(s, Inches(10.0), Inches(1.5), Inches(3.0), Inches(5.5), RGBColor(0xF8,0xFA,0xFC), RGBColor(0xE2,0xE8,0xF0))
tf = set_text(tech_box, "", 11, C_DARKGRAY)
p = tf.paragraphs[0]
p.text = "技术选型"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = C_DARK
p.alignment = PP_ALIGN.CENTER

techs = [
    "前端: React + TypeScript",
    "网关: Go (Gin)",
    "AI层: Python + LangChain",
    "模型: Qwen2.5 / DeepSeek",
    "推理: vLLM 加速",
    "查询: Trino",
    "计算: Spark",
    "表格式: Iceberg",
    "调度: Airflow",
    "向量库: Milvus",
    "元数据: PostgreSQL",
    "缓存: Redis",
    "部署: Kubernetes",
]
for t in techs:
    p = tf.add_paragraph()
    p.text = t
    p.font.size = Pt(10)
    p.font.color.rgb = C_DARKGRAY
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

# ============================================================
# P5 数据目录演进
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "数据目录演进路径", "Data Catalog — 贯穿全平台的数据资产基座，四阶段持续深化")

catalog_items = [
    ("CATALOG L1\n技术元数据目录", "2026 Q2-Q3",
     ["表名、字段名、数据类型", "分区策略与统计信息", "Iceberg Catalog 全量同步", "表间关系（外键、Join模式）"],
     ["表结构浏览", "语义搜索", "元数据API"], C_BLUE),
    ("CATALOG L2\n业务语义目录", "2026 Q4-2027 Q1",
     ["业务术语与字段语义映射", "指标定义与计算口径", "查询热度与使用频率统计", "常用查询模式沉淀"],
     ["自然语言搜表", "指标地图", "使用推荐"], C_PURPLE),
    ("CATALOG L3\n智能治理目录", "2027 Q1-Q2",
     ["字段级血缘关系", "敏感级别自动标注", "数据质量评分与历史趋势", "数据Owner与责任矩阵"],
     ["血缘图谱", "敏感标记", "质量评分"], C_GREEN),
    ("CATALOG L4\n全景数据图谱", "2027 Q2-Q3",
     ["ETL管道与数据流关联映射", "数据新鲜度与更新频率", "端到端链路健康状态", "数据价值评估与成本归属"],
     ["链路追踪", "健康仪表盘", "价值评估"], C_AMBER),
]

for i, (title, period, contents, capabilities, color) in enumerate(catalog_items):
    left = Inches(0.3) + i * Inches(3.2)
    top = Inches(1.8)
    w = Inches(3.0)
    h = Inches(5.2)

    card = add_rounded_rect(s, left, top, w, h, C_WHITE, RGBColor(0xE2,0xE8,0xF0))
    add_rect(s, left, top, w, Inches(0.06), color)

    # 标题
    tf = add_text_box(s, left + Inches(0.15), top + Inches(0.15), w - Inches(0.3), Inches(0.65), title, 13, color, True, PP_ALIGN.CENTER)
    # 时间
    add_text_box(s, left + Inches(0.15), top + Inches(0.85), w - Inches(0.3), Inches(0.3), period, 10, C_GRAY, False, PP_ALIGN.CENTER)

    # 分隔线
    add_rect(s, left + Inches(0.2), top + Inches(1.15), w - Inches(0.4), Inches(0.015), RGBColor(0xE2,0xE8,0xF0))

    # 目录内容
    add_text_box(s, left + Inches(0.15), top + Inches(1.25), w - Inches(0.3), Inches(0.25), "目录内容", 10, C_GRAY, True)
    tf_c = add_text_box(s, left + Inches(0.15), top + Inches(1.5), w - Inches(0.3), Inches(1.8), "", 10, C_DARKGRAY)
    for j, c in enumerate(contents):
        p = tf_c.paragraphs[0] if j == 0 else tf_c.add_paragraph()
        p.text = "· " + c
        p.font.size = Pt(10)
        p.font.color.rgb = C_DARKGRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(3)

    # 分隔线
    add_rect(s, left + Inches(0.2), top + Inches(3.4), w - Inches(0.4), Inches(0.015), RGBColor(0xE2,0xE8,0xF0))

    # 核心能力
    add_text_box(s, left + Inches(0.15), top + Inches(3.5), w - Inches(0.3), Inches(0.25), "核心能力", 10, C_GRAY, True)

    # 标签
    for j, cap in enumerate(capabilities):
        row = j // 2
        col = j % 2
        tag_left = left + Inches(0.15) + col * Inches(1.35)
        tag_top = top + Inches(3.8) + row * Inches(0.4)
        tag = add_rounded_rect(s, tag_left, tag_top, Inches(1.25), Inches(0.3), RGBColor(0xF8,0xFA,0xFC))
        set_text(tag, cap, 9, color, True, PP_ALIGN.CENTER)
        tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 箭头 (between cards)
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + w + Inches(0.02), top + Inches(2.5), Inches(0.14), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_GRAY
        arrow.line.fill.background()

# ============================================================
# P6 第一年路线图
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "第一年演进路线（2026.04 — 2027.03）", "四阶段递进式建设，每阶段在前一阶段基础上展开")

phase_details = [
    ("第一阶段 · 智能SQL开发中心", "2026 Q2-Q3", C_BLUE,
     ["平台基础架构搭建与部署", "Agent框架与工具协议定义", "NL2SQL核心能力开发", "SQL安全审查引擎",
      "📖 技术元数据目录（L1）", "表结构采集 · 语义搜索 · 元数据API"]),
    ("第二阶段 · 智能分析中心", "2026 Q4-2027 Q1", C_PURPLE,
     ["多轮对话式数据探索", "智能可视化推荐引擎", "洞察自动发现引擎", "分析报告自动生成",
      "📖 业务语义目录（L2）", "业务术语映射 · 指标地图 · 自然语言搜表"]),
    ("第三阶段 · 数据治理中心", "2027 Q1-Q2", C_GREEN,
     ["自动化血缘图谱", "数据质量智能监控", "敏感数据智能识别",
      "📖 智能治理目录（L3）", "血缘图谱 · 敏感标记 · 质量评分", "核心数据域治理体系验证"]),
    ("第四阶段 · 管道编排中心", "2027 Q2-Q3", C_AMBER,
     ["自然语言驱动管道设计", "Spark作业 + Airflow DAG自动生成", "智能调度优化引擎",
      "📖 全景数据图谱（L4）", "链路追踪 · 健康仪表盘 · 价值评估", "全链路可观测与智能诊断"]),
]

for i, (title, period, color, items) in enumerate(phase_details):
    left = Inches(0.3) + i * Inches(3.2)
    top = Inches(1.7)
    w = Inches(3.0)

    # 时间轴圆点 + 线
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, left + w/2 - Inches(0.1), Inches(1.55), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    # 横线
    if i < 3:
        line = add_rect(s, left + w/2 + Inches(0.1), Inches(1.63), Inches(3.0 - 0.2), Inches(0.03), RGBColor(0xE2,0xE8,0xF0))

    # 卡片
    card_h = Inches(5.0)
    card = add_rounded_rect(s, left, top + Inches(0.2), w, card_h, C_WHITE, RGBColor(0xE2,0xE8,0xF0))
    # header
    hdr = add_rect(s, left + Inches(0.04), top + Inches(0.24), w - Inches(0.08), Inches(0.55), color)
    set_text(hdr, title, 12, C_WHITE, True, PP_ALIGN.CENTER)
    # period
    add_text_box(s, left + Inches(0.1), top + Inches(0.85), w - Inches(0.2), Inches(0.25), period, 10, C_GRAY, False, PP_ALIGN.CENTER)
    # items
    tf = add_text_box(s, left + Inches(0.15), top + Inches(1.2), w - Inches(0.3), card_h - Inches(1.5), "", 11, C_DARKGRAY)
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "· " + item
        p.font.size = Pt(10)
        p.font.color.rgb = C_DARKGRAY
        if item.startswith("📖"):
            p.font.color.rgb = color
            p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)

# ============================================================
# P7 第二年路线图
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "第二年演进路线（2027.04 — 2028.03）", "完成管道编排建设，进入全面深化与规模化推广阶段")

y2_phases = [
    ("第四阶段 · 管道编排交付", "2027 Q2-Q3", C_AMBER,
     ["自然语言驱动管道设计", "Spark作业 + Airflow DAG自动生成", "智能调度优化引擎", "全链路可观测与智能诊断",
      "📖 全景数据图谱上线", "四大中心完整闭环发布"]),
    ("深化优化 · 智能增强", "2027 Q4", C_INDIGO,
     ["基于用户反馈的模型持续优化", "跨Agent复杂任务协同增强", "知识图谱深化建设", "性能与稳定性全面调优"]),
    ("规模化推广 · 生态扩展", "2028 Q1", RGBColor(0x0E,0xA5,0xE9),
     ["全企业范围推广覆盖", "开放工具生态接入", "行业场景适配与定制", "平台运营指标全面达标"]),
]

for i, (title, period, color, items) in enumerate(y2_phases):
    left = Inches(0.3) + i * Inches(4.2)
    top = Inches(1.7)
    w = Inches(4.0)
    card_h = Inches(5.0)

    # 时间轴圆点
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, left + w/2 - Inches(0.1), Inches(1.55), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    if i < 2:
        line = add_rect(s, left + w/2 + Inches(0.1), Inches(1.63), Inches(4.2 - 0.2), Inches(0.03), RGBColor(0xE2,0xE8,0xF0))

    card = add_rounded_rect(s, left, top + Inches(0.2), w, card_h, C_WHITE, RGBColor(0xE2,0xE8,0xF0))
    hdr = add_rect(s, left + Inches(0.04), top + Inches(0.24), w - Inches(0.08), Inches(0.55), color)
    set_text(hdr, title, 13, C_WHITE, True, PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.1), top + Inches(0.85), w - Inches(0.2), Inches(0.25), period, 10, C_GRAY, False, PP_ALIGN.CENTER)

    tf = add_text_box(s, left + Inches(0.15), top + Inches(1.2), w - Inches(0.3), card_h - Inches(1.5), "", 11, C_DARKGRAY)
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "· " + item
        p.font.size = Pt(11)
        p.font.color.rgb = C_DARKGRAY
        if item.startswith("📖"):
            p.font.color.rgb = color
            p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)

# ============================================================
# P8 关键里程碑
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_WHITE)
add_section_header(s, "关键里程碑总览", "19个关键里程碑，每个对应明确的交付物和验收标准")

milestones = [
    ("M1", "2026.05", "架构搭建完成", "第一阶段", "基础架构就绪，Agent框架、工具协议、前后端骨架可运行"),
    ("M1.5", "2026.06", "技术元数据目录上线", "数据目录", "Iceberg Catalog全量同步，语义搜索、元数据API可用"),
    ("M2", "2026.07", "核心能力可用", "第一阶段", "NL2SQL、SQL执行、安全审查功能上线"),
    ("M3", "2026.08", "灰度试用", "第一阶段", "1-2个业务域灰度试用，收集反馈"),
    ("M4", "2026.09", "SQL中心正式发布", "第一阶段", "覆盖数据团队日常80%查询需求"),
    ("M5", "2026.09", "分析能力上线", "第二阶段", "多轮对话探索、智能可视化功能上线"),
    ("M6", "2026.10", "洞察引擎上线", "第二阶段", "自动洞察发现、分析报告生成功能上线"),
    ("M7", "2026.12", "业务分析试点", "第二阶段", "2-3个业务线完成分析自助化试点"),
    ("M8", "2027.01", "分析中心正式发布", "第二阶段", "业务自助分析占比 ≥ 60%"),
    ("M8.5", "2027.01", "业务语义目录上线", "数据目录", "业务术语映射、指标地图、自然语言搜表"),
    ("M9", "2027.02", "智能治理目录上线", "数据目录", "血缘图谱、敏感标记、质量评分融合至目录"),
    ("M10", "2027.02", "质量与安全上线", "第三阶段", "数据质量监控和敏感数据识别功能上线"),
    ("M11", "2027.03", "治理体系验证", "第三阶段", "核心数据域治理体系完整验证"),
    ("M12", "2027.04", "治理中心正式发布", "第三阶段", "核心数据资产100%入目入册"),
    ("M13", "2027.04", "管道设计上线", "第四阶段", "自然语言驱动管道设计功能上线"),
    ("M13.5", "2027.05", "全景数据图谱上线", "数据目录", "端到端链路追踪、健康仪表盘上线"),
    ("M14", "2027.06", "调度与监控上线", "第四阶段", "智能调度优化和全链路可观测上线"),
    ("M15", "2027.08", "端到端验证", "第四阶段", "从需求到产出的端到端流程验证"),
    ("M16", "2027.09", "平台完整闭环", "第四阶段", "四大中心完整闭环，全面发布"),
]

# 表头
header_top = Inches(1.5)
cols = [Inches(0.5), Inches(1.4), Inches(3.0), Inches(4.8), Inches(6.4)]
col_widths = [Inches(0.9), Inches(1.5), Inches(1.7), Inches(1.5), Inches(4.5)]
headers = ["编号", "时间", "里程碑", "阶段", "核心交付物"]

for j, (x, hdr_text) in enumerate(zip(cols, headers)):
    hdr = add_rect(s, x, header_top, col_widths[j], Inches(0.35), C_DARK)
    set_text(hdr, hdr_text, 10, C_WHITE, True, PP_ALIGN.CENTER)

# 数据行
row_h = Inches(0.27)
tag_colors = {"第一阶段": C_BLUE, "第二阶段": C_PURPLE, "第三阶段": C_GREEN, "第四阶段": C_AMBER, "数据目录": RGBColor(0x0E,0xA5,0xE9)}

for i, (mid, date, name, phase, desc) in enumerate(milestones):
    row_top = header_top + Inches(0.35) + i * row_h
    bg_color = C_WHITE if i % 2 == 0 else RGBColor(0xF8,0xFA,0xFC)

    vals = [mid, date, name, phase, desc]
    for j, (x, val) in enumerate(zip(cols, vals)):
        cell = add_rect(s, x, row_top, col_widths[j], row_h, bg_color)
        color = C_DARKGRAY
        bold = False
        if j == 3:  # 阶段列
            color = tag_colors.get(val, C_DARKGRAY)
            bold = True
        set_text(cell, val, 8.5, color, bold, PP_ALIGN.CENTER if j < 3 else PP_ALIGN.LEFT)

# ============================================================
# P9 KPI
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_LIGHT)
add_section_header(s, "关键成功指标", "可量化、可度量、可考核的年度目标")

kpis = [
    ("NL2SQL 准确率", "≥ 90%", "标准测试集 + 线上抽样评估", 90, C_BLUE),
    ("业务自助分析占比", "≥ 60%", "业务侧自行完成的分析需求占比", 60, C_PURPLE),
    ("核心资产血缘覆盖率", "100%", "纳入血缘追踪的核心数据表占比", 100, C_GREEN),
    ("数据管道开发效率", "5x 提升", "同等复杂度管道的开发工时对比", 80, C_AMBER),
    ("需求交付周期", "天 → 小时", "从需求提出到结果交付的平均时长", 75, C_RED),
    ("用户满意度", "≥ 4.2 / 5.0", "季度满意度调研评分", 84, C_INDIGO),
    ("数据目录资产收录率", "100%", "核心数据资产纳入目录管理的比例", 100, RGBColor(0x0E,0xA5,0xE9)),
    ("自然语言搜表准确率", "≥ 85%", "用户用业务语言搜索到正确数据表的命中率", 85, RGBColor(0x0F,0x17,0x2A)),
]

for i, (label, value, desc, pct, color) in enumerate(kpis):
    row = i // 4
    col = i % 4
    left = Inches(0.3) + col * Inches(3.2)
    top = Inches(1.7) + row * Inches(2.8)
    w = Inches(3.0)
    h = Inches(2.5)

    card = add_rounded_rect(s, left, top, w, h, C_WHITE)
    add_rect(s, left, top, w, Inches(0.06), color)

    add_text_box(s, left + Inches(0.2), top + Inches(0.2), w - Inches(0.4), Inches(0.3), label, 12, C_GRAY)
    add_text_box(s, left + Inches(0.2), top + Inches(0.55), w - Inches(0.4), Inches(0.6), value, 28, color, True)
    add_text_box(s, left + Inches(0.2), top + Inches(1.15), w - Inches(0.4), Inches(0.5), desc, 10, C_GRAY)

    # 进度条
    bar_bg = add_rounded_rect(s, left + Inches(0.2), top + Inches(1.8), w - Inches(0.4), Inches(0.12), RGBColor(0xF0,0xF0,0xF0))
    bar_fill_w = int((w - Inches(0.4)) * pct / 100)
    add_rounded_rect(s, left + Inches(0.2), top + Inches(1.8), bar_fill_w, Inches(0.12), color)

# ============================================================
# P10 结尾
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s, C_DARK)
add_text_box(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1), "DataAgent", 48, C_WHITE, True, PP_ALIGN.CENTER)
add_text_box(s, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6), "让每个人都拥有自己的数据智能伙伴", 20, RGBColor(0x94,0xA3,0xB8), False, PP_ALIGN.CENTER)
add_rect(s, Inches(5.6), Inches(4.3), Inches(2), Inches(0.04), C_BLUE)
add_text_box(s, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4), "2026年4月  |  数据智能平台项目组", 14, RGBColor(0x64,0x74,0x8B), False, PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
output_path = "/home/ouyan/llm/dataagent/docs/DataAgent战略规划蓝图.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
